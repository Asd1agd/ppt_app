from pptx import Presentation
import google.generativeai as genai
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import random



# Configure Gemini
api_key = input("Enter your Gemini API key: ")
genai.configure(api_key="AIzaSyBibUXKKGzdH-mErjMjNgYDY0kuxe4pK_I")

model = genai.GenerativeModel('models/gemini-2.5-flash-preview-04-17')
prompt = input("Enter the PPT title: ")

# Generate content for the presentation
try:
    response = model.generate_content(f"give me presentation agenda list of 6 elements about: {prompt} , and give me the list separated by ',' and noting else not description or else just 10 elements separated by ','")
    content = response.text.strip()
except Exception as e:
    print(f"Error generating content: {e}")
    exit()

# Split the content into agenda items
agenda_items = content.split(",")
print(agenda_items)

# def estimate_text_width(text, font_size_pt):
#     # Estimate width: adjust factor per font (this is rough)
#     average_char_width_in_inches = font_size_pt * 0.0057  # tweak if needed
#     return len(text) * average_char_width_in_inches

# Load template
template = Presentation("formatPpt2.pptx")
new_ppt = Presentation()
new_ppt.slide_width = template.slide_width
new_ppt.slide_height = template.slide_height


def Title(slide,data,font = 'Century Gothic',font_size =54,clr =[0,0,0], Top = 2.5):
    # Fixed left
    left = Inches(6.65)
    top = Inches(Top)
    font_size_pt = font_size
    sp = "                                                                                    \n"
    text = f"{sp}{data}"

    # Add textbox
    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    # Style
    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = f'{font}'
    p.font.bold = True
    p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
    p.alignment = PP_ALIGN.LEFT

    p = text_frame.paragraphs[1]
    p.font.size = Pt(font_size_pt)
    p.font.name = f'{font}'
    p.font.bold = True
    p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
    p.alignment = PP_ALIGN.CENTER

def heading(slide,data,font = 'Century Gothic',font_size =32,clr =[0,0,0], Top = -0.3):
    # Fixed left
    left = Inches(6.65)
    top = Inches(Top)
    font_size_pt = font_size
    sp = "                                                                                    \n"
    text = f"{sp}  {data}"

    # Add textbox
    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    # Style
    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = f'{font}'
    p.font.bold = True
    p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
    p.alignment = PP_ALIGN.LEFT

    p = text_frame.paragraphs[1]
    p.font.size = Pt(font_size_pt)
    p.font.name = f'{font}'
    p.font.bold = True
    p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
    p.alignment = PP_ALIGN.LEFT

    # textbox.fill.solid()
    # textbox.fill.fore_color.rgb = RGBColor(255, 255, 0)

def content(slide,data,font = 'Century Gothic',font_size=20,clr =[0,0,0], Top = 1,space=""):
    # Fixed left
    left = Inches(6.65)
    top = Inches(Top)
    font_size_pt = font_size
    sp = "                                                                                    \n"
    text = f"{sp}{data}"
    text = text.replace("\n",f"\n  {space}")

    # Add textbox
    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    # Style
    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = f'{font}'
    p.font.bold = True
    p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
    p.alignment = PP_ALIGN.LEFT

    for i in range(len(text.split('\n')) - 1):
        p = text_frame.paragraphs[i + 1]
        p.font.size = Pt(font_size_pt)
        p.font.name = f'{font}'
        p.font.bold = True
        p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
        p.alignment = PP_ALIGN.LEFT

    # textbox.fill.solid()
    # textbox.fill.fore_color.rgb = RGBColor(255, 255, 0)

# Title Slide
slide1 = new_ppt.slides.add_slide(template.slides[1].slide_layout)
Title(slide1,prompt)

# Agendas Slide
slide2 = new_ppt.slides.add_slide(template.slides[10].slide_layout)
heading(slide2,"               "+"Agenda")
agenda_content = "\n".join([f"{i+1}. {agenda}" for i, agenda in enumerate(agenda_items)])
content(slide2,agenda_content,space="\n       ")


# Create individual slides for each agenda item
for i, agenda_item in enumerate(agenda_items):
    if agenda_item.strip():  # Skip any empty agenda item
        print(f"[INFO] Generating slide for agenda {i + 1}: {agenda_item}")
        # Create a new slide with heading and content for each agenda item
        slide_no = random.randint(2, 6)
        slide = new_ppt.slides.add_slide(template.slides[slide_no].slide_layout)
        Agenda_no = f"Agenda {i+1}: {agenda_item.strip()}"
        heading(slide, Agenda_no)

        # Generate content for each agenda item (optional, adjust logic as needed)
        agenda_content_response = model.generate_content(f"Provide short content for the agenda item in 5 to 6 points: {agenda_item.strip()} and just give me the content points nothing else not any coment or description from your side")
        agenda_content = agenda_content_response.text.strip()

        # Remove any unwanted formatting (like stars for bold)
        agenda_content = agenda_content.replace("*", "")  # Remove stars used for bold formatting

        content(slide,agenda_content,space="\n    ")

        print(f"Slide {i+1} generated")

# Title Slide
slideL = new_ppt.slides.add_slide(template.slides[-1].slide_layout)
Title(slideL,"Thanks")

# Save the presentation to a file
ppt_io = BytesIO()
new_ppt.save(ppt_io)

with open("generated.pptx", "wb") as f:
    f.write(ppt_io.getbuffer())

print("Presentation generated: generated.pptx")