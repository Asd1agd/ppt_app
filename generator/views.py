from django.shortcuts import render
from django.http import HttpResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import google.generativeai as genai
import random

imgPath = "formatbg.png"

# --- Fix: Pass `prs` to get dimensions from presentation object
def copy_slide(slide, image_path, prs):
    # Use the correct method to get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Add background image
    bg_image = slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)

    # Send image to back
    slide.shapes._spTree.remove(bg_image._element)
    slide.shapes._spTree.insert(2, bg_image._element)  # Put it behind other elements
# kk

# ---- Helper functions for text ----
def ppt_Title(slide, data, font='Century Gothic', font_size=54, clr=[0, 0, 0], Top=2.5):
    left = Inches(6.65)
    top = Inches(Top)
    text = f"{' '*90}\n{data}"

    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.LEFT

    p = text_frame.paragraphs[1]
    p.font.size = Pt(font_size)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.CENTER

def ppt_heading(slide, data, font='Century Gothic', font_size=32, clr=[0, 0, 0], Top=0.25):
    left = Inches(6.65/10)
    top = Inches(Top)
    text = f"{' '*90}\n      {data}"

    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.LEFT

    p = text_frame.paragraphs[1]
    p.font.size = Pt(font_size)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.LEFT

def ppt_content(slide, data, font='Century Gothic', font_size=20, clr=[0, 0, 0], Top=1, space=""):
    left = Inches(6.65/10)
    top = Inches(Top)
    text = f"{' '*90}\n{data}".replace("\n", f"\n  {space}")

    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.LEFT

    for i in range(len(text.split('\n')) - 1):
        p = text_frame.paragraphs[i + 1]
        p.font.size = Pt(font_size)
        p.font.name = font
        p.font.bold = True
        p.font.color.rgb = RGBColor(*clr)
        p.alignment = PP_ALIGN.LEFT



def home(request):
    if request.method == 'POST':
        # Load template
        template = Presentation("formatPpt.pptx")
        new_ppt = Presentation()
        new_ppt.slide_width = template.slide_width
        new_ppt.slide_height = template.slide_height

        api_key = request.POST['api_key']
        api_key = "AIzaSyBibUXKKGzdH-mErjMjNgYDY0kuxe4pK_I"
        title = request.POST['title']

        # Configure Gemini API
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('models/gemini-2.5-flash-preview-04-17')

        # Prompt for the presentation agenda generation
        prompt = (
            f"give me presentation agenda list of 6 elements about: {title}, "
            f"and give me the list separated by ',' and nothing else. "
            f"No description, just 6 elements separated by commas."
        )
        try:
            # Generate content for the agenda
            response = model.generate_content(prompt)
            content = response.text.strip()
            print(f"[INFO] Agenda generated: {content}")
        except Exception as e:
            print(f"[ERROR] Gemini API failed: {e}")
            return HttpResponse(f"<h3>Error generating content: {e}</h3>")

        # Split the content into agenda items
        agenda_items = [item.strip() for item in content.split(",")]
        print(f"[INFO] Agenda items parsed: {agenda_items}")

        # Title Slide
        slide1 = new_ppt.slides.add_slide(new_ppt.slide_layouts[6])
        copy_slide(slide1, imgPath, new_ppt)
        ppt_Title(slide1, title)

        # Agendas Slide
        slide2 = new_ppt.slides.add_slide(new_ppt.slide_layouts[6])
        copy_slide(slide2, imgPath, new_ppt)
        ppt_heading(slide2, "               " + "Agenda")
        agenda_content = "\n".join([f"{i + 1}. {agenda}" for i, agenda in enumerate(agenda_items)])
        ppt_content(slide2, agenda_content, space="\n       ")

        # Individual agenda slides
        for i, item in enumerate(agenda_items):
            if item.strip():
                print(f"[INFO] Generating slide for agenda {i+1}: {item}")
                # Create a new slide with heading and content for each agenda item
                slide_no = random.randint(3, 7)
                slide = new_ppt.slides.add_slide(new_ppt.slide_layouts[6])
                copy_slide(slide, imgPath, new_ppt)
                Agenda_no = f"Agenda {i + 1}: {item.strip()}"
                ppt_heading(slide, Agenda_no)

                try:
                    detail_response = model.generate_content(
                        f"Provide short content for the agenda item in 4 to 5 points: {item.strip()} and just give me the content points nothing else not any coment or description from your side"
                    )
                    agenda_details = detail_response.text.strip().replace("*", "")
                except Exception as e:
                    agenda_details = "Content generation failed."
                    print(f"[WARN] Failed content for {item}: {e}")

                # Add content to slide
                ppt_content(slide, agenda_details, space="\n       ")

        # Thanks Slide
        slideL = new_ppt.slides.add_slide(new_ppt.slide_layouts[6])
        copy_slide(slideL, imgPath, new_ppt)
        ppt_Title(slideL, "Thanks")

        # Save and return presentation
        ppt_io = BytesIO()
        new_ppt.save(ppt_io)
        ppt_io.seek(0)
        print("[INFO] PPT generation complete")

        response = HttpResponse(
            ppt_io.getvalue(),
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response['Content-Disposition'] = 'attachment; filename=generated.pptx'
        print("[INFO] PPT sent for download")
        return response

    return render(request, 'form.html')
