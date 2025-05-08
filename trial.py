# AIzaSyBibUXKKGzdH-mErjMjNgYDY0kuxe4pK_I

import google.generativeai as genai
from pptx import Presentation
from pptx.util import Pt
from io import BytesIO

# Configure Gemini
api_key = input("Enter your Gemini API key: ")
genai.configure(api_key="AIzaSyBibUXKKGzdH-mErjMjNgYDY0kuxe4pK_I")

model = genai.GenerativeModel('models/gemini-2.5-flash-preview-04-17')
prompt = input("Enter the PPT title: ")

# Generate content for the presentation
try:
    response = model.generate_content(f"give me presentation agenda list of 6 elements about: {prompt} , and give me the list separated by ',' and noting else not description or else just 10 elements separated by ','")
    content = response.text.strip()
except Exception as e:
    print(f"Error generating content: {e}")
    exit()

# Split the content into agenda items
agenda_items = content.split(",")
print(agenda_items)

# Create a PowerPoint presentation
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = prompt
slide.placeholders[1].text = "An AI-generated presentation"

# Agendas Slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Agenda"
agenda_content = "\n".join([f"{i+1}. {agenda}" for i, agenda in enumerate(agenda_items)])
slide.placeholders[1].text = agenda_content

# Helper function to set font size and fit text
def set_font_size(placeholder, font_size=18):
    # Accessing the text frame directly
    if placeholder.has_text_frame:
        for paragraph in placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)

# Create individual slides for each agenda item
for i, agenda_item in enumerate(agenda_items):
    if agenda_item.strip():  # Skip any empty agenda item
        # Create a new slide with heading and content for each agenda item
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Agenda {i+1}: {agenda_item.strip()}"

        # Generate content for each agenda item (optional, adjust logic as needed)
        agenda_content_response = model.generate_content(f"Provide short content for the agenda item in 3 to 4 points: {agenda_item.strip()} and just give me the content points nothing else not any coment or description from your side")
        agenda_content = agenda_content_response.text.strip()

        # Remove any unwanted formatting (like stars for bold)
        agenda_content = agenda_content.replace("*", "")  # Remove stars used for bold formatting

        # Add the content related to the agenda item
        textbox = slide.shapes.placeholders[1]
        textbox.text = agenda_content

        # Set font size for the content
        set_font_size(textbox, font_size=18)

        # Ensure content fits within the slide
        if len(agenda_content) > 800:  # If content is too long, trim it
            agenda_content = agenda_content[:800] + "..."
            textbox.text = agenda_content
            set_font_size(textbox, font_size=18)

        print(f"Slide {i+1} generated")

# Save the presentation to a file
ppt_io = BytesIO()
prs.save(ppt_io)

with open("generated.pptx", "wb") as f:
    f.write(ppt_io.getbuffer())

print("Presentation generated: generated.pptx")






















#
#
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# from pptx.enum.text import PP_ALIGN
#
# def estimate_text_width(text, font_size_pt):
#     # Estimate width: adjust factor per font (this is rough)
#     average_char_width_in_inches = font_size_pt * 0.0057  # tweak if needed
#     return len(text) * average_char_width_in_inches
#
# # Load template
# template = Presentation("formatPpt2.pptx")
# new_ppt = Presentation()
# new_ppt.slide_width = template.slide_width
# new_ppt.slide_height = template.slide_height
#
# # Add slide
# slide2 = new_ppt.slides.add_slide(template.slides[2].slide_layout)
#
# def heading(slide,data,font = 'Century Gothic',font_size =32,clr =[0,0,0], Top = -0.3):
#     # Fixed left
#     left = Inches(6.65)
#     top = Inches(Top)
#     font_size_pt = font_size
#     sp = "                                                                                    \n"
#     text = f"{sp}  {data}"
#
#     # Estimate required width
#     estimated_width_in = estimate_text_width(text, font_size_pt)
#     width = Inches(estimated_width_in)
#     height = Inches(1)
#
#     # Add textbox
#     textbox = slide.shapes.add_textbox(left, top, 0, 0)
#     text_frame = textbox.text_frame
#     text_frame.text = text
#
#     # Style
#     p = text_frame.paragraphs[0]
#     p.font.size = Pt(32)
#     p.font.name = f'{font}'
#     p.font.bold = True
#     p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
#     p.alignment = PP_ALIGN.LEFT
#
#     p = text_frame.paragraphs[1]
#     p.font.size = Pt(font_size_pt)
#     p.font.name = f'{font}'
#     p.font.bold = True
#     p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
#     p.alignment = PP_ALIGN.LEFT
#
#     # textbox.fill.solid()
#     # textbox.fill.fore_color.rgb = RGBColor(255, 255, 0)
#
# def content(slide,data,font = 'Century Gothic',font_size=20,clr =[0,0,0], Top = 1):
#     # Fixed left
#     left = Inches(6.65)
#     top = Inches(Top)
#     font_size_pt = font_size
#     sp = "                                                                                    \n"
#     text = f"{sp}{data}"
#     text = text.replace("\n","\n  ")
#
#     # Estimate required width
#     estimated_width_in = estimate_text_width(text, font_size_pt)
#     width = Inches(estimated_width_in)
#     height = Inches(1)
#
#     # Add textbox
#     textbox = slide.shapes.add_textbox(left, top, 0, 0)
#     text_frame = textbox.text_frame
#     text_frame.text = text
#
#     # Style
#     p = text_frame.paragraphs[0]
#     p.font.size = Pt(32)
#     p.font.name = f'{font}'
#     p.font.bold = True
#     p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
#     p.alignment = PP_ALIGN.LEFT
#
#     for i in range(len(data.split('\n'))):
#         p = text_frame.paragraphs[i + 1]
#         p.font.size = Pt(font_size_pt)
#         p.font.name = f'{font}'
#         p.font.bold = True
#         p.font.color.rgb = RGBColor(clr[0], clr[1], clr[2])
#         p.alignment = PP_ALIGN.LEFT
#
#     # textbox.fill.solid()
#     # textbox.fill.fore_color.rgb = RGBColor(255, 255, 0)
#
# heading(slide2,"Title")
#
# content(slide2,"This is para 1.\ni am the\ncoder.")
#
#
# # Save
# new_ppt.save("new_presentation.pptx")

