from pptx import Presentation
import google.generativeai as genai
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
import random

imgPath = "formatbg.png"

# --- Fix: Pass `prs` to get dimensions from presentation object
def copy_slide(slide, image_path, prs):
    # Use the correct method to get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Add background image
    bg_image = slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)

    # Send image to back
    slide.shapes._spTree.remove(bg_image._element)
    slide.shapes._spTree.insert(2, bg_image._element)  # Put it behind other elements

# Configure Gemini
api_key = input("Enter your Gemini API key: ")
genai.configure(api_key="AIzaSyBibUXKKGzdH-mErjMjNgYDY0kuxe4pK_I")

model = genai.GenerativeModel('models/gemini-2.5-flash-preview-04-17')
prompt = input("Enter the PPT title: ")

# Generate agenda content
try:
    response = model.generate_content(f"give me presentation agenda list of 6 elements about: {prompt} , and give me the list separated by ',' and noting else not description or else just 6 elements separated by ','")
    content = response.text.strip()
except Exception as e:
    print(f"Error generating content: {e}")
    exit()

agenda_items = content.split(",")
print(agenda_items)

# Load template and new presentation
template = Presentation("formatPpt.pptx")
new_ppt = Presentation()
new_ppt.slide_width = template.slide_width
new_ppt.slide_height = template.slide_height

# ---- Helper functions for text ----
def Title(slide, data, font='Century Gothic', font_size=54, clr=[0, 0, 0], Top=2.5):
    left = Inches(6.65)
    top = Inches(Top)
    text = f"{' '*90}\n{data}"

    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.LEFT

    p = text_frame.paragraphs[1]
    p.font.size = Pt(font_size)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.CENTER

def heading(slide, data, font='Century Gothic', font_size=32, clr=[0, 0, 0], Top=0.25):
    left = Inches(6.65)
    top = Inches(Top)
    text = f"{' '*90}\n      {data}"

    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.LEFT

    p = text_frame.paragraphs[1]
    p.font.size = Pt(font_size)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.LEFT

def content(slide, data, font='Century Gothic', font_size=20, clr=[0, 0, 0], Top=1, space=""):
    left = Inches(6.65)
    top = Inches(Top)
    text = f"{' '*90}\n{data}".replace("\n", f"\n  {space}")

    textbox = slide.shapes.add_textbox(left, top, 0, 0)
    text_frame = textbox.text_frame
    text_frame.text = text

    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = font
    p.font.bold = True
    p.font.color.rgb = RGBColor(*clr)
    p.alignment = PP_ALIGN.LEFT

    for i in range(len(text.split('\n')) - 1):
        p = text_frame.paragraphs[i + 1]
        p.font.size = Pt(font_size)
        p.font.name = font
        p.font.bold = True
        p.font.color.rgb = RGBColor(*clr)
        p.alignment = PP_ALIGN.LEFT

# Title Slide
slide1 = new_ppt.slides.add_slide(new_ppt.slide_layouts[6])
copy_slide(slide1, imgPath, new_ppt)
Title(slide1, prompt)

# Agenda Slide
slide2 = new_ppt.slides.add_slide(new_ppt.slide_layouts[6])
copy_slide(slide2, imgPath, new_ppt)
heading(slide2, "Agenda")
agenda_content = "\n".join([f"{i + 1}. {agenda.strip()}" for i, agenda in enumerate(agenda_items)])
content(slide2, agenda_content, space="\n       ")

# Agenda Item Slides
for i, agenda_item in enumerate(agenda_items):
    if agenda_item.strip():
        print(f"[INFO] Generating slide for agenda {i + 1}: {agenda_item}")
        slide = new_ppt.slides.add_slide(new_ppt.slide_layouts[6])
        copy_slide(slide, imgPath, new_ppt)

        heading(slide, f"Agenda {i + 1}: {agenda_item.strip()}")

        try:
            agenda_content_response = model.generate_content(
                f"Provide short content for the agenda item in 5 to 6 points: {agenda_item.strip()} and just give me the content points nothing else"
            )
            agenda_content = agenda_content_response.text.strip().replace("*", "")
        except Exception as e:
            agenda_content = "[Failed to generate content]"

        content(slide, agenda_content, space="\n    ")
        print(f"[INFO] Slide {i + 1} generated")

# Closing Slide
slideL = new_ppt.slides.add_slide(new_ppt.slide_layouts[6])
copy_slide(slideL, imgPath, new_ppt)
Title(slideL, "Thanks")

# Save
ppt_io = BytesIO()
new_ppt.save(ppt_io)

with open("generated.pptx", "wb") as f:
    f.write(ppt_io.getbuffer())

print("✅ Presentation generated: generated.pptx")
